"""
Text extraction service for document files.
"""
import logging
import os
from typing import Dict, Optional, Tuple

import PyPDF2
from pptx import Presentation

# Setup logging
logger = logging.getLogger(__name__)


class ExtractionService:
    """Service for extracting text from document files."""

    def extract_text(self, file_path: str) -> Tuple[Optional[str], Optional[Dict[str, str]]]:
        """Extract text from a document file.

        Args:
            file_path: Path to the document file.

        Returns:
            Tuple of (full text, page/slide content dictionary).
        """
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None, None

        file_extension = os.path.splitext(file_path)[1].lower()

        try:
            if file_extension == ".pdf":
                return self._extract_from_pdf(file_path)
            elif file_extension in [".pptx", ".ppt"]:
                return self._extract_from_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_extension}")
                return None, None
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return None, None

    def _extract_from_pdf(self, pdf_path: str) -> Tuple[str, Dict[str, str]]:
        """Extract text from a PDF file.

        Args:
            pdf_path: Path to the PDF file.

        Returns:
            Tuple of (full text, page content dictionary).
        """
        logger.info(f"Extracting text from PDF: {pdf_path}")
        full_text = ""
        page_content = {}

        try:
            with open(pdf_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                num_pages = len(reader.pages)

                for page_num in range(num_pages):
                    page = reader.pages[page_num]
                    page_text = page.extract_text() or ""

                    # Store page content
                    page_content[str(page_num + 1)] = page_text

                    # Add to full text
                    full_text += page_text + "\n\n"

            logger.info(f"Successfully extracted text from {num_pages} pages")
            return full_text.strip(), page_content
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            raise

    def _extract_from_pptx(self, pptx_path: str) -> Tuple[str, Dict[str, str]]:
        """Extract text from a PowerPoint file.

        Args:
            pptx_path: Path to the PowerPoint file.

        Returns:
            Tuple of (full text, slide content dictionary).
        """
        logger.info(f"Extracting text from PowerPoint: {pptx_path}")
        full_text = ""
        slide_content = {}

        try:
            presentation = Presentation(pptx_path)

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""

                # Extract text from all shapes that contain text
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                # Store slide content
                slide_content[str(slide_num)] = slide_text.strip()

                # Add to full text
                full_text += slide_text + "\n\n"

            logger.info(f"Successfully extracted text from {len(presentation.slides)} slides")
            return full_text.strip(), slide_content
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {str(e)}")
            raise
