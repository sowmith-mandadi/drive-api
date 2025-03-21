"""PowerPoint content extractor module."""

import os
import logging
from typing import List, Dict, Any
from pptx import Presentation

# Initialize logger
logger = logging.getLogger(__name__)

class PptxExtractor:
    """Extracts content from PowerPoint files."""
    
    def __init__(self):
        """Initialize the PowerPoint extractor."""
        logger.info("Initializing PowerPoint Extractor")
    
    def extract_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text content from a PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List[Dict]: List of text chunks with slide information
        """
        try:
            logger.info(f"Processing PowerPoint file: {file_path}")
            
            if not os.path.exists(file_path):
                logger.error(f"File not found: {file_path}")
                return []
            
            # Create PowerPoint presentation object
            presentation = Presentation(file_path)
            
            # Extract text from each slide
            chunks = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                
                # Combine all text from the slide
                if slide_text:
                    combined_text = "\n".join(slide_text)
                    
                    # Skip empty slides
                    if not combined_text.strip():
                        continue
                    
                    # Add to chunks
                    chunks.append({
                        "text": combined_text,
                        "slide": i + 1
                    })
            
            logger.info(f"Extracted {len(chunks)} text chunks from PowerPoint")
            return chunks
        
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {e}")
            return [] 